'''
    powerpoint example
    https://www.geeksforgeeks.org/creating-and-updating-powerpoint-presentations-in-python-using-python-pptx/
    https://towardsdatascience.com/creating-presentations-with-python-3f5737824f61
    https://pbpython.com/creating-powerpoint.html
    https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
    Ref for slide types:
        0 ->  title and subtitle
        1 ->  title and content
        2 ->  section header
        3 ->  two content
        4 ->  Comparison
        5 ->  Title only
        6 ->  Blank
        7 ->  Content with caption
        8 ->  Pic with caption
'''
from pptx import Presentation


def main():
    '''
        main function
    '''
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Hello, World!"
    slide.placeholders[1].text = "python-pptx was here!"
    prs.save('test.pptx')


if __name__ == "__main__":
    main()
